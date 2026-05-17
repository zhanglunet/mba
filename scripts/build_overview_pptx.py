#!/usr/bin/env python3
"""Build MBA project overview deck as .pptx.

Focuses on the 5 areas:
1. Problem MBA solves
2. System architecture
3. How the program was developed
4. Sample report (Lenovo)
5. How to use
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "hackathon" / "mba-overview-deck.pptx"
SCREENSHOTS = ROOT / "docs" / "screenshots"
IMAGES = ROOT / "images"

# Colors
BG          = RGBColor(0x0B, 0x0D, 0x10)
PANEL       = RGBColor(0x14, 0x18, 0x1E)
INK         = RGBColor(0xE9, 0xEE, 0xF5)
INK_DIM     = RGBColor(0x9A, 0xA6, 0xB2)
LINE        = RGBColor(0x1F, 0x26, 0x30)
ACCENT      = RGBColor(0xFF, 0x7A, 0x4D)
ACCENT2     = RGBColor(0x67, 0xE8, 0xA4)
WARN        = RGBColor(0xF1, 0xC4, 0x53)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_SANS = "PingFang SC"
FONT_MONO = "Menlo"


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=PANEL, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *,
             font=FONT_SANS, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if "\n" in text else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


def add_runs(slide, x, y, w, h, runs, *,
             font=FONT_SANS, size=18, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.3):
    """runs = list of (text, color, bold) — all in one paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, color, bold in runs:
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


def add_chrome(slide, page_num, total, tag=None):
    """Page number, footer, optional tag."""
    add_text(slide, Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.3),
             f"{page_num:02d} / {total:02d}", font=FONT_MONO, size=10,
             color=INK_DIM, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "MBA · Metric Brand Auditor", font=FONT_MONO, size=10, color=INK_DIM)
    if tag:
        add_text(slide, Inches(0.5), Inches(0.35), Inches(6), Inches(0.3),
                 tag.upper(), font=FONT_MONO, size=11, color=ACCENT)


def add_accent_bar(slide, x=Inches(0.5), y=Inches(0.85), h=Inches(0.06), w=Inches(0.6)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_h1(slide, text, y=Inches(1.0)):
    add_text(slide, Inches(0.5), y, Inches(12.3), Inches(1.2), text,
             size=44, bold=True, color=INK, line_spacing=1.15)


def add_h2(slide, text, y=Inches(1.0)):
    add_text(slide, Inches(0.5), y, Inches(12.3), Inches(0.9), text,
             size=32, bold=True, color=INK, line_spacing=1.2)


def add_section_divider(prs, title, subtitle, num):
    """Big section transition slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_text(s, Inches(0.5), Inches(2.5), Inches(3), Inches(0.6),
             f"SECTION 0{num}", font=FONT_MONO, size=14, color=ACCENT)
    # big section title
    add_text(s, Inches(0.5), Inches(3.1), Inches(12.3), Inches(1.5),
             title, size=64, bold=True, color=INK, line_spacing=1.1)
    add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.8),
             subtitle, size=22, color=INK_DIM, line_spacing=1.4)
    # accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.4),
                              Inches(0.8), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    return s


def add_simple_table(slide, x, y, w, h, headers, rows, *, col_widths=None,
                     header_color=INK_DIM, row_color=INK, accent_first=False):
    """Hand-rolled table. col_widths sums to w, in EMU."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    if col_widths is None:
        col_widths = [w // n_cols] * n_cols
    row_h = h // n_rows
    cur_x = x
    # header row
    for i, label in enumerate(headers):
        add_text(slide, cur_x, y, col_widths[i], row_h, label,
                 size=12, color=header_color, font=FONT_MONO, bold=True)
        cur_x += col_widths[i]
    # underline
    line = slide.shapes.add_connector(1, x, y + row_h - Inches(0.05),
                                      x + w, y + row_h - Inches(0.05))
    line.line.color.rgb = LINE
    line.line.width = Pt(0.75)
    # body rows
    for ri, row in enumerate(rows):
        cur_x = x
        for ci, cell in enumerate(row):
            if isinstance(cell, tuple):
                # (text, color, bold)
                txt, col, bd = cell
                add_text(slide, cur_x, y + (ri + 1) * row_h, col_widths[ci], row_h,
                         txt, size=15, color=col, bold=bd, line_spacing=1.3)
            else:
                col = ACCENT if accent_first and ci == 0 else row_color
                bd  = accent_first and ci == 0
                add_text(slide, cur_x, y + (ri + 1) * row_h, col_widths[ci], row_h,
                         cell, size=15, color=col, bold=bd, line_spacing=1.3)
            cur_x += col_widths[ci]
        # row separator
        ly = y + (ri + 2) * row_h - Inches(0.05)
        if ri < len(rows) - 1:
            sep = slide.shapes.add_connector(1, x, ly, x + w, ly)
            sep.line.color.rgb = LINE
            sep.line.width = Pt(0.25)


def add_code_block(slide, x, y, w, h, code):
    box = add_rect(slide, x, y, w, h, fill=RGBColor(0x0A, 0x0C, 0x10), line=LINE)
    pad = Inches(0.25)
    add_text(slide, x + pad, y + pad, w - 2 * pad, h - 2 * pad, code,
             font=FONT_MONO, size=14, color=RGBColor(0xD7, 0xE3, 0xF4), line_spacing=1.45)


def add_bullets(slide, x, y, w, h, items, *, size=18, color=INK, bullet="•",
                bullet_color=ACCENT, line_spacing=1.5):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        # bullet run
        rb = p.add_run()
        rb.text = f"{bullet}  "
        rb.font.name = FONT_MONO
        rb.font.size = Pt(size)
        rb.font.color.rgb = bullet_color
        rb.font.bold = True
        # content runs
        if isinstance(it, list):
            for txt, col, bd in it:
                r = p.add_run()
                r.text = txt
                r.font.name = FONT_SANS
                r.font.size = Pt(size)
                r.font.color.rgb = col
                r.font.bold = bd
        else:
            r = p.add_run()
            r.text = it
            r.font.name = FONT_SANS
            r.font.size = Pt(size)
            r.font.color.rgb = color


# ============================================================================
# Build deck
# ============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

TOTAL = 26  # update if slide count changes

# ---------- 1. Cover ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
# accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.6),
                         Inches(0.1), Inches(2.3))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
add_text(s, Inches(0.8), Inches(2.5), Inches(8), Inches(0.5),
         "PROJECT OVERVIEW · 2026", font=FONT_MONO, size=14, color=ACCENT)
add_runs(s, Inches(0.8), Inches(3.0), Inches(12), Inches(1.5),
         [("MBA", INK, True), (" — ", INK_DIM, False),
          ("Metric Brand Auditor", ACCENT, True)], size=60, line_spacing=1.15)
add_text(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.8),
         "把品牌影响力做成可调研、可打分、可比较、可复盘的智能资产",
         size=22, color=INK, line_spacing=1.4)
add_text(s, Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.5),
         "Audit pipeline that turns brand influence into a versioned, scoreable signal asset",
         size=14, color=INK_DIM)
# meta strip
meta_y = Inches(6.4)
meta = [
    ("Site",       "mbabrand.com"),
    ("BotLearn",   "metric-brand-auditor"),
    ("Team",       "Jason · 清风 · John"),
    ("Co-pilot",   "Claude Opus 4.7"),
]
for i, (k, v) in enumerate(meta):
    x = Inches(0.8 + i * 3.0)
    add_text(s, x, meta_y, Inches(2.8), Inches(0.3), k.upper(),
             font=FONT_MONO, size=10, color=INK_DIM)
    add_text(s, x, meta_y + Inches(0.3), Inches(2.8), Inches(0.4), v,
             font=FONT_SANS, size=14, color=INK, bold=True)

# ---------- 2. TOC ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 2, TOTAL, tag="Contents")
add_accent_bar(s)
add_h1(s, "本场分 5 个部分讲完")
sections = [
    ("01", "要解决的问题",     "Brand influence is unauditable today",         ACCENT),
    ("02", "系统架构",         "5-phase pipeline · 4-layer repo",              ACCENT),
    ("03", "程序如何开发的",   "Multi-agent + perspective DNA + governance",   ACCENT),
    ("04", "样例报告",         "Lenovo 0992.HK · 5-judge audit",               ACCENT),
    ("05", "怎么使用",         "One-line CLI · zero-dep entry · evolution",    ACCENT),
]
y0 = Inches(2.3)
for i, (num, title, sub, c) in enumerate(sections):
    y = y0 + Inches(i * 0.85)
    add_text(s, Inches(0.6), y, Inches(0.9), Inches(0.7), num,
             font=FONT_MONO, size=28, color=c, bold=True)
    add_text(s, Inches(1.6), y, Inches(5), Inches(0.7), title,
             size=26, color=INK, bold=True)
    add_text(s, Inches(6.5), y, Inches(6.3), Inches(0.7), sub,
             size=18, color=INK_DIM, line_spacing=1.5)

# ============================================================================
# SECTION 01 — 要解决的问题
# ============================================================================
add_section_divider(prs, "要解决的问题", "Brand influence is important — yet remains unauditable", 1)

# ---------- 4. The 3 existing methods are broken ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 4, TOTAL, tag="Problem · 现状")
add_accent_bar(s)
add_h2(s, "今天大家怎么判断「品牌」？三种方式都有硬伤")

add_simple_table(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.6),
    headers=["现有方式", "看起来很有用", "但回答不了"],
    rows=[
        ["咨询公司报告",
         "深度访谈、行业对比、漂亮 deck",
         "一次性 · 主观 · 不可复盘 · 半年后没法核对"],
        ["舆情 / 社媒监测",
         "声量、曝光、情绪分布",
         "声量高 ≠ 影响力 · 没有品类定义权诊断"],
        ["创始人 / 投资人自评",
         "经验丰富、直觉准",
         "无统一口径 · 不可跨公司比 · 不可跨时间比"],
    ],
    col_widths=[Inches(3.0), Inches(4.5), Inches(4.8)],
    accent_first=True)

add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6),
         "财务有审计、产品有数据分析、销售有 CRM、舆情有 monitoring —— 品牌影响力没有审计协议。",
         size=18, color=WARN, bold=True, line_spacing=1.4)

# ---------- 5. The reframe — Agent world ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 5, TOTAL, tag="Problem · 升维")
add_accent_bar(s)
add_h2(s, "更大的变化：AI 时代，品牌正在被 Agent 阅读")

# two columns
col_w = Inches(5.9)
gap = Inches(0.4)
left_x = Inches(0.5)
right_x = left_x + col_w + gap
col_y = Inches(2.3)
col_h = Inches(3.8)

# left card
add_rect(s, left_x, col_y, col_w, col_h, fill=PANEL, line=LINE)
add_text(s, left_x + Inches(0.3), col_y + Inches(0.25), col_w, Inches(0.4),
         "过去", font=FONT_MONO, size=14, color=INK_DIM)
add_text(s, left_x + Inches(0.3), col_y + Inches(0.7), col_w - Inches(0.5), Inches(3),
         "品牌主要影响人。\n\n消费者怎么感知它，投资人怎么理解它，媒体怎么叙述它，候选人愿不愿意加入它。",
         size=18, color=INK, line_spacing=1.55)

# right card
add_rect(s, right_x, col_y, col_w, col_h, fill=PANEL, line=ACCENT)
add_text(s, right_x + Inches(0.3), col_y + Inches(0.25), col_w, Inches(0.4),
         "未来", font=FONT_MONO, size=14, color=ACCENT)
add_text(s, right_x + Inches(0.3), col_y + Inches(0.7), col_w - Inches(0.5), Inches(3),
         "越来越多的判断会被 Agent 参与甚至代理：\n\n投资 Agent 筛公司 · 采购 Agent 选供应商\n招聘 Agent 理解雇主 · 战略 Agent 追踪竞争",
         size=18, color=INK, line_spacing=1.55)

add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8),
         "品牌正在从「人类心智资产」变成「智能体世界里可识别、可比较、可调用的信号资产」。",
         size=20, color=ACCENT, bold=True, line_spacing=1.4)

# ---------- 6. Three problems → three mechanisms ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 6, TOTAL, tag="Problem · 三机制")
add_accent_bar(s)
add_h2(s, "三个老问题 → MBA 的三个机制")

add_simple_table(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4),
    headers=["传统报告的问题", "MBA 的应对"],
    rows=[
        ["单线程单视角",
         "N 路并行 sub-agent + 5 位人物评委独立打分，Lead 只做合成"],
        ["不可复盘",
         "版本化目录 reports/<brand>/versions/v{n}_{date}.{md,html}，evolution 滚动"],
        ["打分主观、口径漂移",
         "固定 5 镜头 × 7 维度，跨品牌、跨评委、跨时间同口径"],
    ],
    col_widths=[Inches(4.5), Inches(7.8)],
    accent_first=True)

# ============================================================================
# SECTION 02 — 系统架构
# ============================================================================
add_section_divider(prs, "系统架构", "5-phase orchestrated pipeline · 4-layer repo · 7×5 scoring grid", 2)

# ---------- 8. 5-phase pipeline ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 8, TOTAL, tag="Architecture · Pipeline")
add_accent_bar(s)
add_h2(s, "流水线 5 阶段（FRESH 模式）")

phases = [
    ("Phase 0", "Router",          "检查 report.md 是否存在 → 走 FRESH 或 EVOLUTION"),
    ("Phase 1", "Discovery",       "Lead 起草 PRD，用户确认维度（GATE 1）"),
    ("Phase 2", "Parallel Search", "一条消息派发 N 个 sub-agent + 1 个 wuying 云浏览器 agent"),
    ("Phase 3", "Synthesis",       "Lead 合成 _raw/，产出 synthesis.md（杠杆地图 + 脆弱边缘 + 矛盾点）"),
    ("Phase 4", "5-Judge Panel",   "并行派发 5 评委 skill，独立打分互不可见，落到 reviews/"),
    ("Phase 5", "Lead Merge",      "产出 report.md + report.html + versions/v{n}.{md,html}"),
]
y0 = Inches(2.25)
for i, (ph, name, desc) in enumerate(phases):
    y = y0 + Inches(i * 0.72)
    # phase number card
    add_rect(s, Inches(0.5), y, Inches(1.4), Inches(0.6), fill=PANEL, line=LINE)
    add_text(s, Inches(0.5), y, Inches(1.4), Inches(0.6), ph,
             font=FONT_MONO, size=14, color=ACCENT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # name
    add_text(s, Inches(2.1), y, Inches(2.5), Inches(0.6), name,
             size=18, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # desc
    add_text(s, Inches(4.7), y, Inches(8.2), Inches(0.6), desc,
             size=15, color=INK_DIM, anchor=MSO_ANCHOR.MIDDLE)

# ---------- 9. 7 dimensions × 5 lenses ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 9, TOTAL, tag="Architecture · Scoring grid")
add_accent_bar(s)
add_h2(s, "评分坐标系：7 维度 × 5 镜头")

# left: 7 dimensions
add_text(s, Inches(0.5), Inches(2.2), Inches(6), Inches(0.4),
         "7 个调研维度（子 agent 跑的横向）", font=FONT_MONO, size=13, color=ACCENT)
dims = [
    "1. 创始 & 起源叙事",
    "2. 产品 & 定位",
    "3. 分发 & 渠道",
    "4. 社区 & PR",
    "5. 视觉 & 语言",
    "6. 竞品 & 格局",
    "7. 接收 & 情绪",
]
for i, d in enumerate(dims):
    add_text(s, Inches(0.5), Inches(2.7 + i * 0.5), Inches(6), Inches(0.5),
             d, size=17, color=INK, line_spacing=1.3)

# right: 5 lenses
add_text(s, Inches(7), Inches(2.2), Inches(6), Inches(0.4),
         "5 个打分镜头（评委的纵向）", font=FONT_MONO, size=13, color=ACCENT2)
lenses = [
    ("01", "原创性",       "Origin authenticity"),
    ("02", "范畴命名",     "Category coinage"),
    ("03", "杠杆质量",     "Leverage quality"),
    ("04", "身份一致性",   "Identity coherence"),
    ("05", "真实信号",     "Real-world signal"),
]
for i, (n, zh, en) in enumerate(lenses):
    y = Inches(2.7 + i * 0.7)
    add_text(s, Inches(7), y, Inches(0.5), Inches(0.6), n,
             font=FONT_MONO, size=14, color=ACCENT2, bold=True)
    add_text(s, Inches(7.5), y, Inches(2.5), Inches(0.6), zh,
             size=18, color=INK, bold=True)
    add_text(s, Inches(10), y, Inches(3), Inches(0.6), en,
             size=14, color=INK_DIM, line_spacing=1.2)

add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
         "维度是「调研的输入」，镜头是「评委的尺子」 — 通过 _raw/synthesis.md 耦合。",
         size=15, color=INK_DIM, font=FONT_MONO)

# ---------- 10. 4-layer repo ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 10, TOTAL, tag="Architecture · Repo")
add_accent_bar(s)
add_h2(s, "仓库 4 层架构")
code = """mba/
├── metric-brand-auditor/    ← 编排层：流水线主 SKILL
│   ├── SKILL.md                  Lead 工作手册（Phase 0–5，~600 行）
│   ├── references/               维度模板 / 评委模板 / 云浏览器规范 / HTML 脚手架
│   └── reports/<brand>/          每个品牌一个目录
│
├── research/                ← 工具层：PRD 多代理深度调研 skill
│
├── *-perspective/           ← 评委层：5 套人物视角 skill
│   └── fusheng · jobs · likejia · wu-jundong · zhang-yiming
│       每套 = SKILL.md + references/research/01–06.md + scripts/
│
└── wuying_open.py           ← 基建层：阿里云无影 AgentBay 云浏览器"""
add_code_block(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(4.4), code)
add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
         "每一层对应流水线一个阶段，且可独立调用（research / perspective 都能单独 / 调）。",
         size=15, color=INK_DIM, font=FONT_MONO)

# ---------- 11. 5-judge panel ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 11, TOTAL, tag="Architecture · Judges")
add_accent_bar(s)
add_h2(s, "5 位人物评委")

judges = [
    ("傅盛",      "猎豹 / OpenClaw",      "AI 应用为王 + 战略三板斧",          "傅盛.jpg"),
    ("Steve Jobs","Apple / NeXT / Pixar", "Liberal Arts × Tech · focus = no", "jobs.jpg"),
    ("李可佳",    "BotLearn / Aibrary",   "协议位 + 新物种命名",               "李可佳.jpg"),
    ("吴俊东",    "Aibrary 联创",         "教育本质 + 长期关系投资",           "吴俊东.jpg"),
    ("张一鸣",    "字节跳动",             "Context not Control · 火星视角",   "张一鸣.jpg"),
]

card_w = Inches(2.36)
card_h = Inches(4.1)
gap = Inches(0.1)
x0 = Inches(0.5)
y0 = Inches(2.1)
for i, (name, org, tag, img) in enumerate(judges):
    x = x0 + (card_w + gap) * i
    add_rect(s, x, y0, card_w, card_h, fill=PANEL, line=LINE)
    img_path = IMAGES / img
    if img_path.exists():
        try:
            s.shapes.add_picture(str(img_path), x + Inches(0.3), y0 + Inches(0.3),
                                 width=card_w - Inches(0.6), height=Inches(1.9))
        except Exception:
            pass
    add_text(s, x, y0 + Inches(2.35), card_w, Inches(0.5), name,
             size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(2.85), card_w, Inches(0.4), org,
             size=12, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_MONO)
    add_text(s, x + Inches(0.2), y0 + Inches(3.3), card_w - Inches(0.4), Inches(0.8),
             tag, size=12, color=INK_DIM, align=PP_ALIGN.CENTER, line_spacing=1.4)

add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
         "每套 skill 基于 80%+ 一手访谈与文章；SKILL.md 顶部明确 anti-fabrication 红线 — 不替本人编造未公开内容。",
         size=14, color=INK_DIM, font=FONT_MONO, line_spacing=1.5)

# ---------- 12. Outputs ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 12, TOTAL, tag="Architecture · Outputs")
add_accent_bar(s)
add_h2(s, "每次运行产出 5 类文件")

add_simple_table(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.2),
    headers=["文件", "作用"],
    rows=[
        [("report.md", ACCENT2, True),                          "当前 canonical Markdown 报告（滚动覆盖）"],
        [("report.html", ACCENT2, True),                        "自包含 HTML：Chart.js 雷达 + 异议热力图 + Mermaid 流程图"],
        [("versions/v{n}_{date}.{md,html}", ACCENT2, True),     "每次 evolution 的不可变快照（时间序列）"],
        [("reviews/<judge>.md", ACCENT2, True),                 "5 位评委的独立打分卡"],
        [("_raw/dimension_n_*.md", ACCENT2, True),              "每个维度子 agent 的原始输出（可审计）"],
    ],
    col_widths=[Inches(4.5), Inches(7.8)])

# ============================================================================
# SECTION 03 — 程序如何开发的
# ============================================================================
add_section_divider(prs, "程序如何开发的", "Skill-based · multi-agent · 80% primary sources · governance built-in", 3)

# ---------- 14. Build philosophy ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 14, TOTAL, tag="Build · 开发哲学")
add_accent_bar(s)
add_h2(s, "三条不可妥协的开发原则")

cards = [
    ("01", "Skill, not script",
     "每一层都是 Claude Code skill — SKILL.md 是工作手册，子 agent 通过 SKILL 加载世界观。可以独立调用、可以彼此 LOAD。"),
    ("02", "并行 over 串行",
     "Phase 2 一条消息派发 N 个 sub-agent，Phase 4 一条消息派发 5 评委。Lead 不读中间结果，只合成最终 _raw/。"),
    ("03", "一手 over 二手",
     "5 套 perspective 的 references/research/01–06.md 80%+ 是一手访谈 / 文章 / 播客 transcript，附 anti-fabrication 红线。"),
]
y0 = Inches(2.2)
card_h = Inches(1.5)
for i, (num, title, desc) in enumerate(cards):
    y = y0 + Inches(i * 1.65)
    add_rect(s, Inches(0.5), y, Inches(12.3), card_h, fill=PANEL, line=LINE)
    add_text(s, Inches(0.8), y + Inches(0.2), Inches(0.8), Inches(0.5),
             num, font=FONT_MONO, size=22, color=ACCENT, bold=True)
    add_text(s, Inches(1.8), y + Inches(0.2), Inches(10), Inches(0.5),
             title, size=20, color=INK, bold=True)
    add_text(s, Inches(1.8), y + Inches(0.75), Inches(10.3), Inches(0.7),
             desc, size=14, color=INK_DIM, line_spacing=1.55)

# ---------- 15. How each perspective skill was built ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 15, TOTAL, tag="Build · Perspective DNA")
add_accent_bar(s)
add_h2(s, "一套 perspective skill 是怎么炼成的")

# left: 6-route research
add_text(s, Inches(0.5), Inches(2.2), Inches(6), Inches(0.4),
         "6 路并行调研（research skill 自调用）", font=FONT_MONO, size=13, color=ACCENT)
routes = [
    "01-writings.md         本人公开文章 / 长文 / 博客",
    "02-conversations.md    访谈 / 播客 transcript",
    "03-expression-dna.md   口头禅 / 表达风格 / 反共识陈述",
    "04-external-views.md   外部评价 / 媒体侧写 / 同行视角",
    "05-decisions.md        重大决策史 / 转折点 / 公开反思",
    "06-timeline.md         身份切换轨迹 / 关键时间线",
]
for i, r in enumerate(routes):
    add_text(s, Inches(0.5), Inches(2.65 + i * 0.45), Inches(6), Inches(0.4),
             r, font=FONT_MONO, size=12, color=INK, line_spacing=1.3)

# right: SKILL.md structure
add_text(s, Inches(7), Inches(2.2), Inches(6), Inches(0.4),
         "落到 SKILL.md 后包含 4 件事", font=FONT_MONO, size=13, color=ACCENT2)
parts = [
    ("Frontmatter",       "name · description · 显式 / 主题型触发规则"),
    ("Mental models",     "5–6 个核心心智模型（如「协议位」「赛马机制」）"),
    ("Decision heuristics","8–10 条决策启发式 + 反共识陈述"),
    ("Anti-fabrication",  "哪些话题留白 / 不替本人编造 / 与兄弟 skill 的路由"),
]
for i, (k, v) in enumerate(parts):
    y = Inches(2.65 + i * 0.85)
    add_text(s, Inches(7), y, Inches(2.5), Inches(0.4), k,
             size=14, color=INK, bold=True)
    add_text(s, Inches(7), y + Inches(0.35), Inches(5.8), Inches(0.5), v,
             size=12, color=INK_DIM, line_spacing=1.4)

# ---------- 16. The audit loop in code ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 16, TOTAL, tag="Build · 多 agent 调度")
add_accent_bar(s)
add_h2(s, "Lead → sub-agent → judge 的协作循环")

code = """# Phase 2 — 并行调研（伪代码）
for dim in active_dimensions:                    # 7 维度中用户确认的子集
    spawn_subagent(
        type="general-purpose",
        prompt=dimensions[dim].prompt_template,   # 来自 references/dimensions.md
        output=f"_raw/dimension_{dim}.md",
    )                                             # 一条消息发起所有 sub-agent

# Phase 4 — 并行打分
for judge in ["fusheng", "jobs", "likejia", "wu-jundong", "zhang-yiming"]:
    spawn_subagent(
        load_skill=f"{judge}-perspective",         # 加载该评委的世界观
        prompt=judge_prompt_template.format(synthesis_path="_raw/synthesis.md"),
        output=f"reviews/{judge}.md",
    )                                             # 5 评委互不可见，独立打分"""
add_code_block(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(4.6), code)

# ---------- 17. Governance + anti-fabrication ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 17, TOTAL, tag="Build · 治理")
add_accent_bar(s)
add_h2(s, "可治理输出：Anti-fabrication 红线写在 SKILL 顶部")

add_bullets(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4),
            [
                "评委对自己留白的话题（家庭、未公开决策细节）保持沉默",
                "5 套 perspective 都基于公开一手资料 — 访谈、文章、播客 transcript",
                "网站明确声明：评委头像、评分、辩论是 AI 模拟，不代表本人真实意见",
                "不构成投资建议或真实品牌评价 — 演示用途，治理边界外置",
                "商业化版本会去掉真人模拟风险 → 升级为 Investor / Founder / Product / Category / Distribution / Culture Lens",
            ],
            size=18, line_spacing=1.75)

# ============================================================================
# SECTION 04 — 样例报告
# ============================================================================
add_section_divider(prs, "样例报告", "Lenovo 0992.HK · 22.0/50 · 5 judges, σ ≤ 0.63 = 高度共识", 4)

# ---------- 19. Lenovo screenshot ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 19, TOTAL, tag="Sample · Lenovo")
add_accent_bar(s)
add_h2(s, "Lenovo 0992.HK · 实际报告长这样")
img = SCREENSHOTS / "report-lenovo.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.5), Inches(2.1),
                         width=Inches(8.5))
# right column annotations
right_x = Inches(9.3)
add_text(s, right_x, Inches(2.2), Inches(3.6), Inches(0.4),
         "报告含 5 类组件", font=FONT_MONO, size=12, color=ACCENT)
parts = [
    ("雷达图",          "5 镜头评分"),
    ("异议热力图",      "5 评委 × 5 镜头 σ"),
    ("Mermaid 流程图",  "杠杆 → 哪里来 / 流到哪里"),
    ("执行摘要",        "杠杆 / 脆弱 / 矛盾"),
    ("90 天建议",       "可执行 action list"),
]
for i, (k, v) in enumerate(parts):
    y = Inches(2.7 + i * 0.6)
    add_text(s, right_x, y, Inches(2.0), Inches(0.4), k,
             size=14, color=INK, bold=True)
    add_text(s, right_x, y + Inches(0.3), Inches(3.6), Inches(0.4), v,
             size=12, color=INK_DIM)

# ---------- 20. Lenovo score matrix ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 20, TOTAL, tag="Sample · 评分矩阵")
add_accent_bar(s)
add_h2(s, "5 评委 × 5 镜头评分矩阵")

add_simple_table(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(3.8),
    headers=["Lens", "🐆 傅盛", "🍎 Jobs", "🦞 李可佳", "📚 吴俊东", "🐘 张一鸣", "Mean", "σ"],
    rows=[
        ["Origin authenticity", "6", "6", "5", "6", "7", ("6.0", ACCENT, True), "0.63"],
        ["Category coinage",    "3", "3", "2", "3", "3", ("2.8", WARN, True),   "0.40"],
        ["Leverage quality",    "5", "6", "4", "5", "5", ("5.0", ACCENT, True), "0.63"],
        ["Identity coherence",  "3", "3", "3", "3", "4", ("3.2", WARN, True),   "0.40"],
        ["Real-world signal",   "6", "5", "5", "4", "5", ("5.0", ACCENT, True), "0.63"],
        [("Total /50", INK, True),
         ("23", INK, True), ("23", INK, True), ("19", INK, True),
         ("21", INK, True), ("24", INK, True),
         ("22.0", ACCENT2, True), ("1.79", INK_DIM, False)],
    ],
    col_widths=[Inches(3.3), Inches(1.1), Inches(1.1), Inches(1.2),
                Inches(1.2), Inches(1.4), Inches(1.5), Inches(1.5)])

add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
         "σ ≤ 0.63 = 5 评委高度共识 — 分歧不在打分上，而在 Round 2 辩论里浮现的 5 种独立诊断语言。",
         size=14, color=INK_DIM, font=FONT_MONO, line_spacing=1.5)

# ---------- 21. Lenovo key findings ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 21, TOTAL, tag="Sample · 核心结论")
add_accent_bar(s)
add_h2(s, "联想审计：4 条结构性结论")

findings = [
    ("01",
     "全球 PC 龙头 + 子品牌资产撑起的渠道机器",
     "FY24/25 营收 $69.1B · 全球 PC 出货 #1（25.5% 份额，领先 HP 5.7pct）· 承重资产是 ThinkPad（IT 行业仅次于 Apple logo），主品牌只有 2015 年红色矩形。"),
    ("02",
     "品牌叙事比财务面落后一个时代",
     "slogan 14 年 4 换没一个撑过 5 年；「AI PC」无 IDC/Gartner traction；「Hybrid AI」借自 Qualcomm 2023 白皮书 — 最大声 evangelist 但不是 term owner。"),
    ("03",
     "两边政治定价都为负，但资本市场未折价",
     "中文舆论质疑「不够中国」，Pentagon J-2 + FY2025 NDAA 质疑「太中国」— 但 19 家券商 16 Buy / 0 Sell，12 月 +30%。资本和舆论脱钩。"),
    ("04",
     "ISG +63% 是搭顺风车，不是赢市场份额",
     "服务器市场全球同期 +61%，联想份额仍 ~4%（Dell 10% / Supermicro 9.5% / 浪潮 4.1% 真赢家）。「AI PC」定义权在 Microsoft + Intel + Qualcomm。"),
]
y0 = Inches(2.2)
for i, (num, title, body) in enumerate(findings):
    y = y0 + Inches(i * 1.15)
    add_text(s, Inches(0.5), y, Inches(0.8), Inches(0.4), num,
             font=FONT_MONO, size=18, color=ACCENT, bold=True)
    add_text(s, Inches(1.4), y, Inches(11.4), Inches(0.4), title,
             size=17, color=INK, bold=True)
    add_text(s, Inches(1.4), y + Inches(0.45), Inches(11.4), Inches(0.65), body,
             size=12, color=INK_DIM, line_spacing=1.5)

# ============================================================================
# SECTION 05 — 怎么使用
# ============================================================================
add_section_divider(prs, "怎么使用", "One line. Zero deps. Evolves over time.", 5)

# ---------- 23. Commands ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 23, TOTAL, tag="Usage · 命令")
add_accent_bar(s)
add_h2(s, "一行命令调用")

code = """/mba <brand>                # 标准全流程（自动判断 FRESH / EVOLUTION）
/mba OpenClaw               # 仓库内置 demo case
/mba <brand> --quick        # 跳过云浏览器 leg（只走开放网）
/mba <brand> --refresh      # 强制 EVOLUTION 重跑（已有报告归档到 versions/）
/mba <brand> --no-judges    # 只做合成，跳过 5 评委
/mba <brand> --focus 1,3,7  # 只调研指定维度
/mba list                   # 列出已审计品牌 + 各自版本数"""
add_code_block(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(3.3), code)

# zero-dep entry box
add_rect(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0),
         fill=PANEL, line=ACCENT)
add_text(s, Inches(0.8), Inches(5.95), Inches(3), Inches(0.4),
         "ZERO-DEP ENTRY", font=FONT_MONO, size=11, color=ACCENT, bold=True)
add_text(s, Inches(0.8), Inches(6.25), Inches(11.5), Inches(0.5),
         "/mba <brand> --quick --no-judges  ← 装完直接跑，只用 WebSearch + WebFetch",
         font=FONT_MONO, size=15, color=INK)

# ---------- 24. FRESH vs EVOLUTION ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 24, TOTAL, tag="Usage · 两种模式")
add_accent_bar(s)
add_h2(s, "两种模式：第一次跑 vs 后续追踪")

col_w = Inches(5.9)
left_x  = Inches(0.5)
right_x = Inches(6.9)
y0 = Inches(2.2)
ch = Inches(4.5)

# FRESH
add_rect(s, left_x, y0, col_w, ch, fill=PANEL, line=LINE)
add_text(s, left_x + Inches(0.3), y0 + Inches(0.2), col_w, Inches(0.4),
         "FRESH MODE", font=FONT_MONO, size=14, color=ACCENT, bold=True)
add_text(s, left_x + Inches(0.3), y0 + Inches(0.65), col_w, Inches(0.5),
         "第一次审计这个品牌", size=20, color=INK, bold=True)
fresh_items = [
    "Phase 0 路由器找不到 report.md",
    "Phase 1 起草 PRD，用户确认 7 维度",
    "Phase 2 并行调研所有维度",
    "Phase 3 Lead 合成 synthesis.md",
    "Phase 4 5 评委独立打分",
    "Phase 5 输出 v1 + canonical report",
]
for i, it in enumerate(fresh_items):
    add_text(s, left_x + Inches(0.5), y0 + Inches(1.4 + i * 0.45),
             col_w - Inches(0.6), Inches(0.4), f"·  {it}",
             size=14, color=INK_DIM, line_spacing=1.4)

# EVOLUTION
add_rect(s, right_x, y0, col_w, ch, fill=PANEL, line=ACCENT2)
add_text(s, right_x + Inches(0.3), y0 + Inches(0.2), col_w, Inches(0.4),
         "EVOLUTION MODE", font=FONT_MONO, size=14, color=ACCENT2, bold=True)
add_text(s, right_x + Inches(0.3), y0 + Inches(0.65), col_w, Inches(0.5),
         "后续追踪品牌变化", size=20, color=INK, bold=True)
evo_items = [
    "Phase 0 发现 report.md 已存在",
    "Phase 1E 列 diff plan：哪些维度可能变了",
    "只重跑变了的维度（省时间省成本）",
    "评委只在受影响维度上重打分",
    "Phase 5 版本号 +1，写 v{n}_<date>.{md,html}",
    "时间序列：跟踪品牌如何演化",
]
for i, it in enumerate(evo_items):
    add_text(s, right_x + Inches(0.5), y0 + Inches(1.4 + i * 0.45),
             col_w - Inches(0.6), Inches(0.4), f"·  {it}",
             size=14, color=INK_DIM, line_spacing=1.4)

# ---------- 25. Target users & scenarios ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 25, TOTAL, tag="Usage · 用户场景")
add_accent_bar(s)
add_h2(s, "给谁用、什么时候用")

add_simple_table(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4),
    headers=["用户", "典型场景"],
    rows=[
        ["创始人 / CEO Office",     "战略叙事一致性诊断 / 融资 deck 前的品牌镜子"],
        ["投资人 / VC / PE",        "Brand Influence Due Diligence / portfolio 持续监控"],
        ["战略部 / IR",             "上市公司品牌叙事审计 / 竞品定位追踪 / IPO Narrative Check"],
        ["品牌 / 增长团队",         "发布前复盘 / 品类定义权验证 / Founder Story Audit"],
        ["咨询公司 / 二级研究员",   "输入素材层，加速尽调与战略项目"],
    ],
    col_widths=[Inches(4.5), Inches(7.8)],
    accent_first=True)

# ---------- 26. Closing ----------
s = prs.slides.add_slide(blank); set_slide_bg(s)
add_chrome(s, 26, TOTAL, tag="Thank You")
# big tagline
add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1),
         "我们不是在做 AI 品牌报告。",
         size=36, color=INK_DIM, line_spacing=1.3)
add_text(s, Inches(0.5), Inches(3.1), Inches(12.3), Inches(1.4),
         "我们是在做 AI 时代的品牌影响力审计协议。",
         size=42, color=ACCENT, bold=True, line_spacing=1.2)

# accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.6),
                         Inches(0.8), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

# meta links
links = [
    ("Site",        "mbabrand.com"),
    ("Sample",      "mbabrand.com/reports/lenovo"),
    ("Install",     "botlearn.ai/skillhunt/v2/s/metric-brand-auditor"),
    ("Team",        "Jason · 清风 · John"),
    ("Powered by",  "marsdata.ai"),
]
for i, (k, v) in enumerate(links):
    y = Inches(5.0 + i * 0.4)
    add_text(s, Inches(0.5), y, Inches(2), Inches(0.4), k.upper(),
             font=FONT_MONO, size=11, color=INK_DIM)
    add_text(s, Inches(2.2), y, Inches(10), Inches(0.4), v,
             size=14, color=INK, bold=True, font=FONT_MONO)

# ---------- Save ----------
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"OK · wrote {OUT}")
print(f"   slides = {len(prs.slides)}  size = {OUT.stat().st_size / 1024:.1f} KB")
