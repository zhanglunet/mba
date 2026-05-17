#!/usr/bin/env python3
"""Build the MBA hackathon pitch deck (.pptx).

Output: docs/hackathon/mba-hackathon-deck.pptx
Usage:  python3 scripts/build_hackathon_deck.py
Deps:   pip3 install --user python-pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Colors --------------------------------------------------------------------
NAVY = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT = RGBColor(0xE8, 0x7A, 0x3E)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BG = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

JUDGE_COLORS = {
    "fusheng": RGBColor(0xD9, 0x77, 0x06),
    "jobs": RGBColor(0x94, 0xA3, 0xB8),
    "likejia": RGBColor(0xDC, 0x26, 0x26),
    "wu-jundong": RGBColor(0x16, 0xA3, 0x4A),
    "zhang-yiming": RGBColor(0x1E, 0x29, 0x3B),
}

# 16:9 slide size in EMU (1 inch = 914400 EMU)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ROOT = Path(__file__).resolve().parent.parent


def add_filled_rect(slide, x, y, w, h, fill_color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(
    slide, x, y, w, h, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
    font="PingFang SC",
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_multi_text(slide, x, y, w, h, lines, *, size=14, color=INK, font="PingFang SC"):
    """lines: list of (text, bold, indent_level) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, bold, indent = item
        else:
            text, bold, indent = item, False, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = indent
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
        p.space_after = Pt(6)
    return tb


def slide_header(slide, title, subtitle=None):
    """Title bar at top of every content slide."""
    add_filled_rect(slide, 0, 0, SLIDE_W, Inches(0.85), NAVY)
    add_text(
        slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
        title, size=24, bold=True, color=WHITE,
    )
    if subtitle:
        add_text(
            slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.3),
            subtitle, size=12, color=RGBColor(0xCB, 0xD5, 0xE1),
        )


def page_footer(slide, page, total):
    add_text(
        slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
        f"MBA · {page}/{total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_01_title(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)

    # Big mark
    add_text(
        s, Inches(0.6), Inches(2.2), Inches(12), Inches(1.2),
        "MBA", size=88, bold=True, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.6),
        "Metric Brand Auditor",
        size=32, color=ACCENT,
    )
    add_text(
        s, Inches(0.6), Inches(4.2), Inches(12), Inches(1.2),
        "给品牌叙事做体检的 5 评委 AI 团",
        size=22, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(4.7), Inches(12), Inches(1.2),
        "The 5-judge AI panel that audits brand narratives",
        size=14, color=RGBColor(0xCB, 0xD5, 0xE1),
    )

    # bottom strip
    add_filled_rect(s, 0, Inches(7.0), SLIDE_W, Inches(0.5), ACCENT)
    add_text(
        s, Inches(0.5), Inches(7.07), Inches(12), Inches(0.4),
        "github.com/zhanglunet/mba    ·    Hackathon Pitch    ·    2026",
        size=12, bold=True, color=WHITE,
    )


def slide_02_pain(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    slide_header(s, "品牌分析报告的 3 大顽疾", "Why brand analysis is broken today")

    # Three pain cards
    card_w = Inches(3.9)
    card_h = Inches(4.4)
    gap = Inches(0.25)
    start_x = Inches(0.65)
    y = Inches(1.5)

    pains = [
        ("01", "单线程 单视角", "One analyst, one bias",
         "一个人写一份报告\n容易陷入自家叙事\n或个人偏好"),
        ("02", "不可复盘", "Not auditable over time",
         "报告是一次性的\n半年后再看\n不知道哪些已过时"),
        ("03", "打分漂移", "Score drift",
         "每个分析师有自己的尺度\n跨品牌、跨时间\n不可比"),
    ]
    for i, (num, cn_title, en_title, body) in enumerate(pains):
        x = start_x + (card_w + gap) * i
        add_filled_rect(s, x, y, card_w, card_h, WHITE)
        add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(2), Inches(0.6),
                 num, size=44, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.3), y + Inches(1.1), card_w - Inches(0.6), Inches(0.5),
                 cn_title, size=22, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(1.7), card_w - Inches(0.6), Inches(0.4),
                 en_title, size=12, color=MUTED)
        add_text(s, x + Inches(0.3), y + Inches(2.5), card_w - Inches(0.6), Inches(2.0),
                 body, size=14, color=INK)

    add_text(
        s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
        "→ 我们用 3 条对应机制解决",
        size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
    )
    page_footer(s, page, total)


def slide_03_insight(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    slide_header(s, "不靠一个分析师,靠 5 个名人评委", "5 in-character AI judges, scoring independently")

    judges = [
        ("🐆", "傅盛", "fusheng", "看穿渠道与杠杆", "猎豹 / OpenClaw"),
        ("🍎", "Steve Jobs", "jobs", "看穿品类命名", "Apple / NeXT / Pixar"),
        ("🦞", "李可佳", "likejia", "看穿新物种叙事", "BotLearn / Aibrary"),
        ("📚", "吴俊东", "wu-jundong", "看穿教育与长期投资", "Aibrary 联创 / 前 TAL"),
        ("🐘", "张一鸣", "zhang-yiming", "看穿组织与赛马", "字节跳动"),
    ]

    card_w = Inches(2.4)
    card_h = Inches(3.6)
    gap = Inches(0.13)
    total_w = card_w * 5 + gap * 4
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(1.5)

    for i, (emoji, name, slug, role, bg) in enumerate(judges):
        x = start_x + (card_w + gap) * i
        add_filled_rect(s, x, y, card_w, card_h, WHITE)
        # color stripe top
        add_filled_rect(s, x, y, card_w, Inches(0.18), JUDGE_COLORS[slug])
        # emoji
        add_text(s, x, y + Inches(0.35), card_w, Inches(1.1),
                 emoji, size=60, align=PP_ALIGN.CENTER)
        # name
        add_text(s, x, y + Inches(1.6), card_w, Inches(0.5),
                 name, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # role line
        add_text(s, x + Inches(0.15), y + Inches(2.25), card_w - Inches(0.3), Inches(0.5),
                 role, size=13, color=ACCENT, align=PP_ALIGN.CENTER)
        # background
        add_text(s, x + Inches(0.15), y + Inches(2.85), card_w - Inches(0.3), Inches(0.5),
                 bg, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    add_text(
        s, Inches(0.5), Inches(5.7), Inches(12), Inches(0.5),
        "5 评委 IN CHARACTER 独立打分,DO NOT read each other",
        size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
        "分歧本身就是产品价值 —— 共识告诉你品牌强在哪,撕扯告诉你哪里有故事",
        size=14, color=MUTED, align=PP_ALIGN.CENTER,
    )
    page_footer(s, page, total)


def slide_04_pipeline(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    slide_header(s, "5 阶段流水线", "Phase 0 → 5: from brand name to versioned report")

    # Pipeline boxes
    phases = [
        ("Phase 0", "Router", "FRESH or\nEVOLUTION?"),
        ("Phase 1", "Discovery", "Lead 起草 PRD\n+ GATE 1"),
        ("Phase 2", "Parallel\nResearch", "N sub-agents\n并行调研"),
        ("Phase 3", "Synthesis", "Lead 合成\nleverage map"),
        ("Phase 4", "5-Judge\nPanel", "评委独立\n打分"),
        ("Phase 5", "Merge", "report.md\n+ report.html"),
    ]

    box_w = Inches(1.85)
    box_h = Inches(2.1)
    gap = Inches(0.15)
    start_x = (SLIDE_W - (box_w * 6 + gap * 5)) / 2
    y = Inches(2.0)

    for i, (phase, name, body) in enumerate(phases):
        x = start_x + (box_w + gap) * i
        # box
        add_filled_rect(s, x, y, box_w, box_h, WHITE)
        add_filled_rect(s, x, y, box_w, Inches(0.45), NAVY)
        add_text(s, x, y + Inches(0.05), box_w, Inches(0.4),
                 phase, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), y + Inches(0.6), box_w - Inches(0.2), Inches(0.7),
                 name, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), y + Inches(1.3), box_w - Inches(0.2), Inches(0.8),
                 body, size=10, color=INK, align=PP_ALIGN.CENTER)

        # arrow
        if i < 5:
            ax = x + box_w + Inches(0.02)
            ay = y + box_h / 2 - Inches(0.1)
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, ay, gap - Inches(0.04), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    # KPIs
    add_text(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.5),
             "典型一份报告:", size=14, color=MUTED)

    kpis = [
        ("⏱  ~25 min", NAVY),
        ("💰  $4-5 / 份", NAVY),
        ("📊  5 评委 × 5 镜头 × 7 维度", NAVY),
        ("📝  Markdown + 自包含 HTML", NAVY),
    ]
    kpi_w = Inches(2.85)
    for i, (text, color) in enumerate(kpis):
        x = Inches(0.6) + i * Inches(3.0)
        add_filled_rect(s, x, Inches(5.85), kpi_w, Inches(0.7), WHITE)
        add_text(s, x, Inches(5.95), kpi_w, Inches(0.5),
                 text, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    page_footer(s, page, total)


def slide_05_score_matrix(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    slide_header(s, "实际产出:Score Matrix + Dissent Heatmap",
                 "How readers skim a 25-min audit in 30 seconds")

    # Left: matrix
    add_text(s, Inches(0.6), Inches(1.1), Inches(6), Inches(0.4),
             "5×5 评委打分矩阵 + σ 列", size=14, bold=True, color=NAVY)

    headers = ["镜头 / Lens", "傅盛", "Jobs", "李可佳", "吴俊东", "张一鸣", "σ"]
    rows = [
        ("Origin authenticity",       7, 6, 8, 7, 6, 0.7),
        ("Category coinage",          10, 4, 9, 7, 5, 2.3),
        ("Leverage quality",          6, 5, 6, 5, 6, 0.5),
        ("Identity coherence",        4, 9, 6, 5, 5, 1.7),
        ("Real-world signal",         8, 5, 9, 7, 7, 1.4),
    ]

    table_x = Inches(0.6)
    table_y = Inches(1.55)
    cell_w = [Inches(1.6), Inches(0.55), Inches(0.55), Inches(0.6),
              Inches(0.6), Inches(0.7), Inches(0.5)]
    cell_h = Inches(0.5)

    # header row
    cur_x = table_x
    for col, w in zip(headers, cell_w):
        add_filled_rect(s, cur_x, table_y, w, cell_h, NAVY)
        add_text(s, cur_x, table_y + Inches(0.13), w, Inches(0.4),
                 col, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cur_x += w

    def score_color(score):
        # 1=red, 5=yellow, 10=green
        if score <= 4:
            return RGBColor(0xFE, 0xCA, 0xCA)  # red-200
        if score <= 7:
            return RGBColor(0xFE, 0xF3, 0xC7)  # amber-200
        return RGBColor(0xBB, 0xF7, 0xD0)  # green-200

    def sigma_color(sigma):
        if sigma <= 0.8:
            return RGBColor(0xE0, 0xE7, 0xFF)
        if sigma <= 1.5:
            return RGBColor(0xC7, 0xD2, 0xFE)
        return RGBColor(0xA5, 0xB4, 0xFC)

    # data rows
    for r_i, row in enumerate(rows):
        row_y = table_y + cell_h * (r_i + 1)
        cur_x = table_x
        # name cell
        add_filled_rect(s, cur_x, row_y, cell_w[0], cell_h, WHITE)
        add_text(s, cur_x + Inches(0.1), row_y + Inches(0.13), cell_w[0] - Inches(0.1),
                 Inches(0.4), row[0], size=10, color=INK)
        cur_x += cell_w[0]
        # score cells
        for j, score in enumerate(row[1:6]):
            add_filled_rect(s, cur_x, row_y, cell_w[j + 1], cell_h, score_color(score))
            add_text(s, cur_x, row_y + Inches(0.13), cell_w[j + 1], Inches(0.4),
                     str(score), size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
            cur_x += cell_w[j + 1]
        # sigma
        sigma = row[6]
        add_filled_rect(s, cur_x, row_y, cell_w[6], cell_h, sigma_color(sigma))
        add_text(s, cur_x, row_y + Inches(0.13), cell_w[6], Inches(0.4),
                 f"{sigma}", size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # Right: explanation
    rx = Inches(7.3)
    add_text(s, rx, Inches(1.1), Inches(5.5), Inches(0.4),
             "怎么读这张图?", size=14, bold=True, color=NAVY)

    explain_lines = [
        ("σ ≤ 0.8  共识区 — 大家都说这事儿成立", False, 0),
        ("σ 0.8-1.5  小分歧 — 评委有不同侧重", False, 0),
        ("σ > 1.5  撕扯区 — 这里有故事,值得深挖", True, 0),
    ]
    add_multi_text(s, rx, Inches(1.55), Inches(5.5), Inches(1.5),
                   explain_lines, size=13)

    add_text(s, rx, Inches(3.4), Inches(5.5), Inches(0.4),
             "这一个矩阵 = 30 秒看完的「高密度信号」", size=14, bold=True, color=ACCENT)

    add_multi_text(s, rx, Inches(3.95), Inches(5.5), Inches(2.0), [
        ("• 按 σ 列排序,直接看哪行有撕扯", False, 0),
        ("• 点开撕扯行,跳到评委卡片看 quote", False, 0),
        ("• 共识告诉你品牌强在哪,撕扯暴露真相", False, 0),
        ("• 跨品牌 / 跨时间使用同一坐标系,可对比", False, 0),
    ], size=13)

    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
             "完整 HTML 报告还含:Radar 雷达、Mermaid 影响力流程、Quadrant 定位、评委卡墙",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
    page_footer(s, page, total)


def slide_06_discipline(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    slide_header(s, "为什么 MBA 不是「5 个 ChatGPT 套娃」",
                 "The 3 design disciplines that make judges actually disagree")

    disciplines = [
        ("01",
         "评委不是 prompt,是 skill bundle",
         "Each judge is a real skill",
         "每位评委 = 一个独立 perspective skill\n含 80%+ 一手调研(访谈 / 文章 / 播客)\nLOAD SKILL.md 后才进入打分,确保 in-character\n不是「装作傅盛」——是真正读完傅盛 30+ 篇一手资料"),
        ("02",
         "DO NOT read each other",
         "Independent scoring, enforced",
         "Phase 4 prompt 显式要求评委独立\n互相看不到对方 review 文件\n避免趋同效应:5 评委变 1 评委 ×5"),
        ("03",
         "Anti-fabrication 红线",
         "Hard rules against making things up",
         "每个 perspective skill 顶部明确禁区\n「不要激活」/「不可编造」/「留白」\nMBA 调用时显式 re-read,适用于此刻\n违反 = bug,不是 feature"),
    ]

    card_w = Inches(3.9)
    card_h = Inches(4.4)
    gap = Inches(0.25)
    start_x = Inches(0.65)
    y = Inches(1.5)

    for i, (num, cn, en, body) in enumerate(disciplines):
        x = start_x + (card_w + gap) * i
        add_filled_rect(s, x, y, card_w, card_h, WHITE)
        add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(2), Inches(0.7),
                 num, size=44, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.3), y + Inches(1.05), card_w - Inches(0.6), Inches(0.5),
                 cn, size=18, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(1.7), card_w - Inches(0.6), Inches(0.4),
                 en, size=11, color=MUTED)
        add_text(s, x + Inches(0.3), y + Inches(2.25), card_w - Inches(0.6), Inches(2.0),
                 body, size=12, color=INK)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.4),
             "→ 让评委的分歧本身,成为产品价值",
             size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    page_footer(s, page, total)


def slide_07_architecture(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    slide_header(s, "技术架构:4 层 + 0 build step",
                 "Markdown is code; ship by editing files")

    # Layer cards stacked
    layers = [
        ("编排层 / Orchestration", "metric-brand-auditor/SKILL.md (~600 行)",
         "Lead 的 6-phase 工作手册;EVOLUTION 模式、HTML 模板、9 块图表都从这"),
        ("工具层 / Tools", "research/SKILL.md",
         "PRD 多代理深度调研,被 MBA 复用,也可独立 /research 触发"),
        ("评委层 / Judges", "5 × *-perspective/",
         "fusheng / jobs / likejia / wu-jundong / zhang-yiming · 每套含 6 路调研 + 决策启发式"),
        ("基建层 / Infra", "scripts/wuying/open.py · .env",
         "阿里云无影 AgentBay 一次性会话(免费 Lite 即可);.env 不入库"),
    ]

    layer_w = Inches(11.5)
    layer_h = Inches(0.95)
    start_y = Inches(1.4)
    layer_colors = [NAVY, RGBColor(0x32, 0x5B, 0x8C), RGBColor(0x4D, 0x7C, 0xB5), RGBColor(0x6F, 0x9D, 0xCF)]

    for i, (cn_en, file_ref, desc) in enumerate(layers):
        y = start_y + (layer_h + Inches(0.15)) * i
        add_filled_rect(s, Inches(0.9), y, layer_w, layer_h, WHITE)
        add_filled_rect(s, Inches(0.9), y, Inches(0.18), layer_h, layer_colors[i])
        add_text(s, Inches(1.25), y + Inches(0.1), Inches(4), Inches(0.4),
                 cn_en, size=14, bold=True, color=NAVY)
        add_text(s, Inches(1.25), y + Inches(0.45), Inches(4), Inches(0.4),
                 file_ref, size=10, color=ACCENT, font="Menlo")
        add_text(s, Inches(5.2), y + Inches(0.25), Inches(7.0), Inches(0.6),
                 desc, size=11, color=INK)

    # Bottom note
    add_filled_rect(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.6), RGBColor(0xFE, 0xF3, 0xC7))
    add_text(s, Inches(1.1), Inches(6.18), Inches(11.0), Inches(0.4),
             "🚫 No compile · 🚫 No build · 🚫 No DB    ✅ 改 markdown 即生效 · 贡献门槛 = 会写 markdown",
             size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
    page_footer(s, page, total)


def slide_08_built(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    slide_header(s, "已经做完的", "What's shipped — as of v0.2.2")

    # Left half: built
    lx = Inches(0.7)
    add_text(s, lx, Inches(1.1), Inches(6), Inches(0.5),
             "✅ Shipped", size=22, bold=True, color=RGBColor(0x16, 0xA3, 0x4A))

    add_multi_text(s, lx, Inches(1.7), Inches(6), Inches(5), [
        ("5 套 perspective skills 完整实现", True, 0),
        ("• 每套 80%+ 一手调研(共 ~150 条来源)", False, 1),
        ("• 6 路并行研究 + 决策启发式 + 表达 DNA", False, 1),
        ("• 都内置 anti-fabrication 红线", False, 1),
        ("MBA 主流水线 SKILL.md(~600 行)", True, 0),
        ("• Phase 0-5 + EVOLUTION 模式", False, 1),
        ("• 自包含 HTML 报告(9 块图表)", False, 1),
        ("/research 复用 skill", True, 0),
        ("自包含 HTML 报告(Chart.js + Mermaid · 9 块)", True, 0),
        ("• 雷达图 / 异议热力图 / 影响力流程图 / 定位象限", False, 1),
        ("完整文档手册 docs/(2020 行)", True, 0),
        ("• PRD · 架构 · pipeline · usage · 安装 · 规范 · 扩展", False, 1),
        ("MCP 化设计文档(732 行,未来形态)", True, 0),
    ], size=12)

    # Right half: numbers
    rx = Inches(7.5)
    add_text(s, rx, Inches(1.1), Inches(5.5), Inches(0.5),
             "📊 数据", size=22, bold=True, color=NAVY)

    stats = [
        ("5", "perspective skills"),
        ("7", "默认调研维度"),
        ("9", "HTML 报告组件"),
        ("3", "已发版本(v0.1 / 0.2 / 0.2.2)"),
        ("2752", "行文档(docs/)"),
        ("$4-5", "单次审计成本"),
    ]
    sw = Inches(2.5)
    sh = Inches(1.3)
    for i, (num, label) in enumerate(stats):
        col = i % 2
        row = i // 2
        x = rx + col * (sw + Inches(0.2))
        y = Inches(1.7) + row * (sh + Inches(0.2))
        add_filled_rect(s, x, y, sw, sh, WHITE)
        add_text(s, x, y + Inches(0.1), sw, Inches(0.7),
                 num, size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.85), sw, Inches(0.4),
                 label, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    page_footer(s, page, total)


def slide_09_roadmap(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    slide_header(s, "Roadmap", "From skill bundle to MCP server to brand library")

    milestones = [
        ("v0.3", "持续打磨", "更多真实品牌验证 · sub-agent prompt 调优 · 评委辩论模式 RFC", NAVY, "now"),
        ("v0.4", "MCP 化", "把 MBA 包成 MCP server,任何 MCP-capable agent(OpenClaw / Hermes / Cursor)都能装", ACCENT, "1-2 周"),
        ("v0.5", "扩展点", "add_judge / add_dimension 让外部贡献者加自定义评委、自定义维度", ACCENT, "1 周"),
        ("v1.0", "公开发行", "Cloudflare Pages 托管 brand library · npm publish mba-mcp-server · 评委 marketplace RFC", RGBColor(0x16, 0xA3, 0x4A), "3 个月"),
    ]

    bar_w = Inches(11.5)
    item_h = Inches(1.0)
    start_y = Inches(1.4)

    for i, (ver, name, desc, color, eta) in enumerate(milestones):
        y = start_y + i * (item_h + Inches(0.15))
        add_filled_rect(s, Inches(0.9), y, bar_w, item_h, WHITE)
        # version chip
        add_filled_rect(s, Inches(0.9), y, Inches(1.3), item_h, color)
        add_text(s, Inches(0.9), y + Inches(0.25), Inches(1.3), Inches(0.5),
                 ver, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # name + desc
        add_text(s, Inches(2.4), y + Inches(0.1), Inches(4), Inches(0.4),
                 name, size=16, bold=True, color=NAVY)
        add_text(s, Inches(2.4), y + Inches(0.5), Inches(8.5), Inches(0.5),
                 desc, size=11, color=INK)
        # ETA
        add_text(s, Inches(11.0), y + Inches(0.3), Inches(1.3), Inches(0.4),
                 eta, size=12, color=MUTED, align=PP_ALIGN.RIGHT)

    add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
             "📄 完整 MCP 设计文档:docs/mcp-server-design.md (732 行,7 部分 + 2 附录)",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
    page_footer(s, page, total)


def slide_10_why_now(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    slide_header(s, "为什么现在做 MBA", "Why now — the AI agent era reshapes brand visibility")

    add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.6),
             "AI agent 时代的品牌 ≠ 传统品牌",
             size=24, bold=True, color=NAVY)

    blocks = [
        ("过去", "品牌 = 用户记得",
         "广告 + 渠道 + 口碑\n人会去主动搜你的名字\n搜索引擎决定曝光"),
        ("现在", "品牌 = AI agent 推荐",
         "AI 帮人选品牌\nagent 读什么资料,品牌就被怎么呈现\nAI-readability 成为新竞争力"),
        ("未来", "品牌 = 被 agent 引用",
         "你的 founder 故事\n你的产品 positioning\n你的 community 真实度\n会被 5 个 AI 评委看见"),
    ]

    block_w = Inches(3.9)
    block_h = Inches(3.5)
    gap = Inches(0.25)
    start_x = Inches(0.65)
    y = Inches(2.3)

    for i, (era, headline, body) in enumerate(blocks):
        x = start_x + (block_w + gap) * i
        add_filled_rect(s, x, y, block_w, block_h, WHITE)
        # era badge
        era_color = [MUTED, ACCENT, NAVY][i]
        add_filled_rect(s, x + Inches(0.3), y + Inches(0.3), Inches(1.2), Inches(0.45), era_color)
        add_text(s, x + Inches(0.3), y + Inches(0.36), Inches(1.2), Inches(0.4),
                 era, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text(s, x + Inches(0.3), y + Inches(1.0), block_w - Inches(0.6), Inches(0.7),
                 headline, size=18, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(1.85), block_w - Inches(0.6), Inches(1.6),
                 body, size=12, color=INK)

    add_text(s, Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
             "MBA = 给品牌做「AI 友好度 + 多视角」体检的第一个工具",
             size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    page_footer(s, page, total)


def slide_11_demo(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    slide_header(s, "Demo / Try it now",
                 "From 0 to your first audit in 5 minutes")

    # Left: code-style steps
    lx = Inches(0.7)
    add_text(s, lx, Inches(1.2), Inches(6), Inches(0.5),
             "5 分钟跑通", size=20, bold=True, color=NAVY)

    add_filled_rect(s, lx, Inches(1.85), Inches(6), Inches(4.0), RGBColor(0x1E, 0x29, 0x3B))

    code_lines = [
        ("# 1. clone", False, 0),
        ("git clone https://github.com/zhanglunet/mba.git ~/mba", False, 0),
        ("", False, 0),
        ("# 2. configure", False, 0),
        ("cd ~/mba && cp .env.example .env", False, 0),
        ("# fill in WUYING_API_KEY=akm-...", False, 0),
        ("", False, 0),
        ("# 3. install", False, 0),
        ("pip3 install --user agentbay", False, 0),
        ("", False, 0),
        ("# 4. link to Claude Code", False, 0),
        ("ln -s ~/mba/metric-brand-auditor ~/.claude/skills/mba", False, 0),
        ("", False, 0),
        ("# 5. run", False, 0),
        ("/mba YourBrand", True, 0),
    ]

    tb = s.shapes.add_textbox(lx + Inches(0.2), Inches(2.0), Inches(5.7), Inches(3.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    for i, (text, bold, _) in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text or " "
        run.font.size = Pt(11)
        run.font.name = "Menlo"
        run.font.bold = bold
        run.font.color.rgb = ACCENT if bold else RGBColor(0xCB, 0xD5, 0xE1)
        p.space_after = Pt(0)

    # Right: links
    rx = Inches(7.3)
    add_text(s, rx, Inches(1.2), Inches(5.5), Inches(0.5),
             "更多入口", size=20, bold=True, color=NAVY)

    links = [
        ("📦  GitHub", "github.com/zhanglunet/mba"),
        ("📖  README", "/README.md  ·  介绍 + 产品结构"),
        ("📚  完整文档", "/docs/  ·  9 篇 + 2750 行"),
        ("🔭  MCP 路线图", "/docs/mcp-server-design.md"),
        ("⚡  现场跑一次", "/mba YourBrand  →  ~25 min 出报告"),
    ]
    for i, (icon_label, url) in enumerate(links):
        y = Inches(1.85) + i * Inches(0.85)
        add_filled_rect(s, rx, y, Inches(5.5), Inches(0.7), WHITE)
        add_text(s, rx + Inches(0.2), y + Inches(0.07), Inches(5.0), Inches(0.4),
                 icon_label, size=14, bold=True, color=NAVY)
        add_text(s, rx + Inches(0.2), y + Inches(0.4), Inches(5.0), Inches(0.3),
                 url, size=10, color=MUTED, font="Menlo")

    add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
             "现场演示:/mba <一个评委你想看的品牌>",
             size=12, color=ACCENT, align=PP_ALIGN.CENTER)
    page_footer(s, page, total)


def slide_12_team(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)

    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
             "Team & Thanks", size=36, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "三人小队 + AI 共创",
             size=16, color=RGBColor(0xCB, 0xD5, 0xE1))

    # Three team member cards
    members = [
        ("💡", "创意", "Ideation", "Jason", RGBColor(0xE8, 0x7A, 0x3E)),
        ("🛠", "实现", "Implementation", "清风", RGBColor(0x16, 0xA3, 0x4A)),
        ("🧭", "顾问", "Advisor", "John", RGBColor(0x6F, 0x9D, 0xCF)),
    ]

    card_w = Inches(3.9)
    card_h = Inches(3.4)
    gap = Inches(0.25)
    start_x = Inches(0.65)
    y = Inches(2.1)

    for i, (emoji, role_cn, role_en, name, color) in enumerate(members):
        x = start_x + (card_w + gap) * i
        add_filled_rect(s, x, y, card_w, card_h, WHITE)
        # Top color stripe
        add_filled_rect(s, x, y, card_w, Inches(0.18), color)
        # Emoji
        add_text(s, x, y + Inches(0.45), card_w, Inches(1.0),
                 emoji, size=56, align=PP_ALIGN.CENTER)
        # Role CN
        add_text(s, x, y + Inches(1.6), card_w, Inches(0.5),
                 role_cn, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        # Role EN
        add_text(s, x, y + Inches(2.05), card_w, Inches(0.4),
                 role_en, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        # Name
        add_text(s, x, y + Inches(2.55), card_w, Inches(0.7),
                 name, size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # AI co-author line
    add_text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
             "🤖 协作 AI:Claude Opus 4.7  ·  Anthropic",
             size=13, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)

    # Footer CTA
    add_filled_rect(s, 0, Inches(6.4), SLIDE_W, Inches(1.1), ACCENT)
    add_text(s, Inches(0.5), Inches(6.55), Inches(12), Inches(0.5),
             "Thank you · 一起让品牌评估有 outsider 视角",
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
             "github.com/zhanglunet/mba",
             size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_title,
        slide_02_pain,
        slide_03_insight,
        slide_04_pipeline,
        slide_05_score_matrix,
        slide_06_discipline,
        slide_07_architecture,
        slide_08_built,
        slide_09_roadmap,
        slide_10_why_now,
        slide_11_demo,
        slide_12_team,
    ]
    total = len(builders)

    for i, fn in enumerate(builders):
        if i == 0:
            fn(prs, total)
        else:
            fn(prs, i + 1, total)

    out = ROOT / "docs" / "hackathon" / "mba-hackathon-deck.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✅ {out}")
    print(f"   {total} slides · 16:9 · 可用 Keynote / PowerPoint / WPS 打开编辑")


if __name__ == "__main__":
    main()
